from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, UnstructuredPowerPointLoader
from langchain_core.documents import Document
from pathlib import Path
import src.config.settings as settings 
from src.utils.logger import logger
from src.exceptions.pdf_exception import PDFException
from src.utils.performance import (
    measure_performance
)

@measure_performance(
    "PDF Loading"
)
def load_pdf(pdf_name: str):
    try:
        pdf_path = settings.PDF_DIRECTORY / pdf_name
        if not pdf_path.exists():
            logger.error(f"PDF not found: {pdf_name}")
            raise PDFException(f"PDF not found: {pdf_name}")
        
        logger.info(f"Loading PDF File: {pdf_name}")
        loader = PyPDFLoader(str(pdf_path))        
        documents = loader.load()
        
        # Normalize source metadata to just the filename for easier filtering
        for doc in documents:
            doc.metadata["source"] = pdf_name
        
        logger.info(f"Loading PDF File: {pdf_name} - Completed")
        return documents
    except PDFException:
        raise
    except Exception as ex:
        logger.exception(f"Failed to load PDF '{pdf_name}': {ex}")
        raise PDFException(f"Failed to load PDF '{pdf_name}': {ex}") from ex


@measure_performance(
    "Document Loading"
)
def load_document(file_name: str):
    """Load document from multiple formats: PDF, Word (.docx, .doc), PowerPoint (.pptx, .ppt)"""
    try:
        file_path = settings.PDF_DIRECTORY / file_name
        if not file_path.exists():
            logger.error(f"Document not found: {file_name}")
            raise PDFException(f"Document not found: {file_name}")
        
        file_extension = file_path.suffix.lower()
        logger.info(f"Loading {file_extension} file: {file_name}")
        
        # Select appropriate loader based on file extension
        if file_extension == '.pdf':
            loader = PyPDFLoader(str(file_path))
            documents = loader.load()
            
        elif file_extension in ['.docx', '.doc']:
            loader = Docx2txtLoader(str(file_path))
            documents = loader.load()
            # Add page numbers for Word documents (treated as single page)
            for doc in documents:
                if 'page' not in doc.metadata:
                    doc.metadata['page'] = 1
                    
        elif file_extension in ['.pptx', '.ppt']:
            try:
                loader = UnstructuredPowerPointLoader(str(file_path))
                documents = loader.load()
                # Add page numbers for PowerPoint (each slide is a page)
                for idx, doc in enumerate(documents, 1):
                    if 'page' not in doc.metadata:
                        doc.metadata['page'] = idx
            except Exception as ppt_error:
                logger.warning(f"UnstructuredPowerPointLoader failed, using fallback: {ppt_error}")
                # Fallback: Read as text using python-pptx
                from pptx import Presentation
                prs = Presentation(str(file_path))
                documents = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    text_content = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content.append(shape.text)
                    slide_text = "\n".join(text_content)
                    if slide_text.strip():
                        doc = Document(
                            page_content=slide_text,
                            metadata={"page": slide_num}
                        )
                        documents.append(doc)
        else:
            raise PDFException(f"Unsupported file format: {file_extension}. Supported formats: .pdf, .docx, .doc, .pptx, .ppt")
        
        # Normalize source metadata to just the filename for easier filtering
        for doc in documents:
            doc.metadata["source"] = file_name
        
        logger.info(f"Loading {file_extension} file: {file_name} - Completed ({len(documents)} pages/sections)")
        return documents
        
    except PDFException:
        raise
    except Exception as ex:
        logger.exception(f"Failed to load document '{file_name}': {ex}")
        raise PDFException(f"Failed to load document '{file_name}': {ex}") from ex
